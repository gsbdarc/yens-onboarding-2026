#!/usr/bin/env python3
"""Build a 2-slide Day 2 agenda deck in the Day 1 Agenda style.

Every number below was measured off the Day 1 deck's own Agenda slide, not
invented. The layout SECTION_TITLE_AND_DESCRIPTION already supplies the grey
right-hand panel, the GSB Research Hub logo and the baseline rule, so this
script only draws what Day 1 draws on top: the serif "Agenda", the item list,
and one small coloured tick per item.

Day 1 flows its items as paragraphs in a single box and the ticks happen to
line up; its paragraph spacing does not reconcile with the measured tick pitch,
so copying those values would break at any other item count. Each item here
gets its own box positioned on its tick instead -- visually identical, and
correct for any N.

    python build_agenda_deck.py <day1-template.pptx> <out.pptx>
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from build_day2_deck import BODY, MUTED, drop_unused, layout, strip_slides, write  # noqa: E402

L_AGENDA = 1              # SECTION_TITLE_AND_DESCRIPTION on master 0
SERIF = "Source Serif 4"  # what Day 1 sets "Agenda" in

# --- geometry, measured from Day 1 slide 2 -----------------------------------
TITLE_POS = (0.342, 2.002, 4.424, 1.621)   # "Agenda"
TITLE_PT = 48
PART_POS = (0.342, 3.62, 4.424, 0.40)      # our addition: which part
ITEM_X, ITEM_W = 5.142, 4.196              # item text boxes
ITEM_PT = 16
TICK_X, TICK_W, TICK_H = 5.000, 0.099, 0.259
PITCH = 0.5655                             # tick-to-tick, in inches
BLOCK_CENTRE = 2.7565                      # the item block is centred here

# tick colours in Day 1's order; cycles beyond five
TICKS = ["8C1515", "006C80", "F1AD13", "7F2D48", "417865"]


def agenda_slide(prs, part_line, items):
    s = prs.slides.add_slide(layout(prs, L_AGENDA))
    drop_unused(s, set())          # Day 1 keeps none of the placeholders

    # "Agenda", big serif, left
    x, y, w, h = TITLE_POS
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(box.text_frame, [("Agenda", True)], size=TITLE_PT, font=SERIF)

    # which part this is -- not in Day 1, added so the two slides differ
    if part_line:
        x, y, w, h = PART_POS
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        write(box.text_frame, [part_line], size=14, font=BODY, color=MUTED)

    # one tick + one text box per item, centred as a block on BLOCK_CENTRE
    n = len(items)
    block_h = (n - 1) * PITCH + TICK_H
    first_y = BLOCK_CENTRE - block_h / 2

    for i, text in enumerate(items):
        y = first_y + i * PITCH

        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(TICK_X), Inches(y),
                                  Inches(TICK_W), Inches(TICK_H))
        tick.fill.solid()
        tick.fill.fore_color.rgb = RGBColor.from_string(TICKS[i % len(TICKS)])
        tick.line.fill.background()
        tick.shadow.inherit = False

        # centre the label on the tick
        lb = s.shapes.add_textbox(Inches(ITEM_X), Inches(y - 0.115),
                                  Inches(ITEM_W), Inches(TICK_H + 0.23))
        lb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(lb.text_frame, [text], size=ITEM_PT, font=BODY)

    return s


def build(template, out):
    prs = Presentation(template)
    strip_slides(prs)

    agenda_slide(prs, "Part 1 · Where your code runs", [
        "What is inside any computer",
        "Running a script",
        "Research computing resources",
        "Profiling",
        "Why profiling matters",
        "Slurm",
        "Hands-on",
    ])

    agenda_slide(prs, "Part 2 · Doing many things at once", [
        "Parallelization",
        "Job arrays",
        "GPUs",
        "Where to get help",
        "Hands-on",
    ])

    prs.save(out)
    return prs


if __name__ == "__main__":
    tpl = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else "day2-agenda.pptx"
    deck = build(tpl, out)
    print(f"wrote {out}: {len(deck.slides)} slides")
