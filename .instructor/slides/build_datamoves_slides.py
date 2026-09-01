#!/usr/bin/env python3
"""Two replacement slides for Day 2 slides 10-12.

Slides 10, 11 and 12 all carried one message: data moves disk -> RAM -> CPU ->
disk, and the disk legs are the slow ones. Slide 11 was a prose restatement of
what slide 10's diagram already draws, so its four step-labels are folded into
a one-line strip under the graphic. Slide 12 stays its own beat -- it is the
consequence, and the bridge into the resources section.

The step strip is deliberately one line and set in ink rather than the muted
caption grey: it labels the diagram rather than commenting on it, and the
multi-line caption was removed from this slide by hand.

    python build_datamoves_slides.py <day1-template.pptx> <graphics-dir> <out.pptx>
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from build_day2_deck import (  # noqa: E402
    BODY, INK, strip_slides, takeaway, title_only,
)

# content box, matching build_day2_deck_v2.picture_slide
BOX_L, BOX_R = 0.34, 9.66
PIC_TOP = 1.15
STRIP_TOP = 4.62
STRIP_H = 0.42

STEPS = "1 · load from disk    ›    2 · into RAM    ›    3 · CPU works    ›    4 · save back to disk"


def data_moves_slide(prs, gif):
    s = title_only(prs, "How your data moves")

    with Image.open(gif) as im:
        iw, ih = im.size
    avail_w = BOX_R - BOX_L
    avail_h = STRIP_TOP - 0.12 - PIC_TOP
    scale = min(avail_w / iw, avail_h / ih)
    w, h = iw * scale, ih * scale
    s.shapes.add_picture(str(gif), Inches(BOX_L + (avail_w - w) / 2), Inches(PIC_TOP),
                         Inches(w), Inches(h))

    # the four legs, absorbed from the old slide 11
    box = s.shapes.add_textbox(Inches(BOX_L), Inches(STRIP_TOP), Inches(avail_w), Inches(STRIP_H))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = STEPS
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.name = BODY
    run.font.color.rgb = INK
    return s


def build(template, gfxdir, out):
    prs = Presentation(template)
    strip_slides(prs)

    data_moves_slide(prs, Path(gfxdir) / "data-moves.gif")

    # verbatim from the deck's current slide 12 -- no gratuitous rewording
    takeaway(prs, "Disk is about a million times farther away than RAM.", [
        "The CPU reaches RAM in nanoseconds; a disk read takes milliseconds. If your data "
        "does not fit in RAM all at once, the script keeps going back to disk mid-computation "
        "— and that, far more often than a slow CPU, is what makes a job crawl.",
    ])

    prs.save(out)
    return prs


if __name__ == "__main__":
    tpl, gfx = sys.argv[1], sys.argv[2]
    out = sys.argv[3] if len(sys.argv) > 3 else "day2-datamoves.pptx"
    deck = build(tpl, gfx, out)
    print(f"wrote {out}: {len(deck.slides)} slides")
