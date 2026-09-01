#!/usr/bin/env python3
"""Build the combined Day 2 lecture deck (Part 1 + Part 2) as a .pptx.

The deck is built *on top of* an export of the Day 1 Google Slides deck, so it
inherits Day 1's real slide masters, layouts, theme fonts and theme colours
rather than a guess at them. Day 1's own slides are stripped; only the design
survives.

    # export the Day 1 deck from Drive as .pptx, then:
    python build_day2_deck.py day1-template.pptx day2-combined.pptx

Upload the result to Drive with contentMimeType set to the .pptx type and Drive
converts it to a native, editable Google Slides deck.

Design constants below were read off the Day 1 deck, not invented:
  accent bar   #006C80   table header #7F2D48   table band #F4F4F4
  fonts        Source Sans Pro (headings) · Source Sans 3 (body) · Source Code Pro
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- design ----

ACCENT = RGBColor(0x00, 0x6C, 0x80)
TBL_HEAD = RGBColor(0x7F, 0x2D, 0x48)
TBL_BAND = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x5F, 0x67, 0x75)

HEAD = "Source Sans Pro"
BODY = "Source Sans 3"
MONO = "Source Code Pro"

# Layout indices on master 0 of the Day 1 template.
L_TITLE = 0
L_TITLE_BODY = 2
L_SECTION = 4
L_TWO_COL = 11
L_TITLE_ONLY = 18
L_MAIN_POINT = 26


def layout(prs, idx):
    return prs.slide_masters[0].slide_layouts[idx]


def strip_slides(prs):
    """Remove every slide but keep masters, layouts and theme."""
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


def drop_unused(slide, keep):
    """Delete placeholders we did not fill, so Slides shows no empty prompts."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx not in keep:
            ph._element.getparent().remove(ph._element)


def write(tf, blocks, size=15.5, font=BODY, color=None, space_after=6):
    """Fill a text frame from (text, bold) pairs or plain strings, one per line."""
    tf.word_wrap = True
    first = True
    for block in blocks:
        text, bold = block if isinstance(block, tuple) else (block, False)
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_after = Pt(space_after)
        if text == "":
            continue
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        if color is not None:
            run.font.color.rgb = color
    return tf


# ------------------------------------------------------------- builders ----


def title_slide(prs, title, date, presenters, affiliation):
    s = prs.slides.add_slide(layout(prs, L_TITLE))
    drop_unused(s, {0, 1, 2})
    t = s.placeholders[0].text_frame
    t.text = title
    t.paragraphs[0].runs[0].font.size = Pt(44)
    s.placeholders[1].text_frame.text = date
    write(s.placeholders[2].text_frame,
          [(presenters, True), (affiliation, False)], size=12, font=HEAD)
    return s


def section(prs, title, sub=None):
    """Section divider: title plus the teal accent bar, exactly as Day 1."""
    s = prs.slides.add_slide(layout(prs, L_SECTION))
    drop_unused(s, {0})
    tf = s.placeholders[0].text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.name = HEAD
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(403250), Emu(2103120),
                             Emu(128016), Emu(566928))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    if sub:
        box = s.shapes.add_textbox(Inches(0.72), Inches(3.05), Inches(8.6), Inches(0.4))
        write(box.text_frame, [sub], size=13, font=BODY, color=MUTED)
    return s


def title_body(prs, title, blocks, size=15.5):
    s = prs.slides.add_slide(layout(prs, L_TITLE_BODY))
    drop_unused(s, {0, 1})
    tf = s.placeholders[0].text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(31.5)
    tf.paragraphs[0].runs[0].font.name = HEAD
    write(s.placeholders[1].text_frame, blocks, size=size)
    return s


def takeaway(prs, headline, body, kicker="THE TAKEAWAY"):
    """The MAIN_POINT slide: small kicker, big line, rule, supporting text."""
    s = prs.slides.add_slide(layout(prs, L_MAIN_POINT))
    drop_unused(s, {0})
    k = s.shapes.add_textbox(Inches(0.75), Inches(0.72), Inches(6.0), Inches(0.28))
    write(k.text_frame, [(kicker, True)], size=10, font=BODY, color=ACCENT)
    tf = s.placeholders[0].text_frame
    tf.text = headline
    tf.paragraphs[0].runs[0].font.size = Pt(27)
    tf.paragraphs[0].runs[0].font.bold = True
    rule = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.75), Inches(3.42),
                                  Inches(2.20), Inches(3.42))
    rule.line.color.rgb = ACCENT
    rule.line.width = Pt(1.5)
    box = s.shapes.add_textbox(Inches(0.75), Inches(3.62), Inches(8.4), Inches(1.2))
    write(box.text_frame, body, size=13)
    return s


def title_only(prs, title, sub=None):
    s = prs.slides.add_slide(layout(prs, L_TITLE_ONLY))
    drop_unused(s, {0})
    tf = s.placeholders[0].text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(31.5)
    tf.paragraphs[0].runs[0].font.name = HEAD
    if sub:
        box = s.shapes.add_textbox(Inches(0.34), Inches(0.95), Inches(9.32), Inches(0.34))
        write(box.text_frame, [sub], size=13, font=BODY, color=MUTED)
    return s


def table_slide(prs, title, headers, rows, note=None, col_widths=None, sub=None):
    s = title_only(prs, title, sub=sub)
    top = Inches(1.42 if sub else 1.20)
    height = Inches(0.34) * (len(rows) + 1)
    shape = s.shapes.add_table(len(rows) + 1, len(headers),
                               Inches(0.34), top, Inches(9.32), height)
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(Inches(9.32) * w / total))

    def cell(r, c, text, bold, fill, color):
        cl = tbl.cell(r, c)
        cl.fill.solid()
        cl.fill.fore_color.rgb = fill
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.margin_left = cl.margin_right = Inches(0.08)
        cl.margin_top = cl.margin_bottom = Inches(0.03)
        tf = cl.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(11.8)
        run.font.bold = bold
        run.font.name = MONO if text.startswith("`") else BODY
        run.text = text.strip("`")
        run.font.color.rgb = color

    for c, h in enumerate(headers):
        cell(0, c, h, True, TBL_HEAD, WHITE)
    for r, row in enumerate(rows, start=1):
        fill = TBL_BAND if r % 2 == 0 else WHITE
        for c, v in enumerate(row):
            cell(r, c, v, c == 0 and len(headers) > 2, fill, INK)

    if note:
        box = s.shapes.add_textbox(Inches(0.34), Inches(4.44), Inches(9.32), Inches(1.10))
        write(box.text_frame, note, size=11.5, color=MUTED, space_after=3)
    return s


def two_col(prs, title, left_head, left, right_head, right, note=None):
    """Used for the checkpoint slides: everyone on the left, extra on the right."""
    s = prs.slides.add_slide(layout(prs, L_TWO_COL))
    drop_unused(s, {0, 1, 2})
    tf = s.placeholders[0].text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(31.5)
    tf.paragraphs[0].runs[0].font.name = HEAD
    write(s.placeholders[1].text_frame, [(left_head, True)] + left, size=13)
    write(s.placeholders[2].text_frame, [(right_head, True)] + right, size=13)
    if note:
        box = s.shapes.add_textbox(Inches(0.34), Inches(4.72), Inches(9.32), Inches(0.6))
        write(box.text_frame, note, size=11.5, color=MUTED)
    return s


# --------------------------------------------------------- native diagrams --


def _box(slide, x, y, w, h, label, caption=None, fill=None, line=ACCENT,
         label_size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or WHITE
    shp.line.color.rgb = line
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(label_size)
    r.font.bold = bold
    r.font.name = HEAD
    r.font.color.rgb = INK
    if caption:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = caption
        r2.font.size = Pt(10)
        r2.font.name = BODY
        r2.font.color.rgb = MUTED
    return shp


def _arrow(slide, x1, y, x2, label=None, color=ACCENT):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y),
                                   Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(1.75)
    if label:
        lw = max(1.9, x2 - x1)
        mid = (x1 + x2) / 2
        box = slide.shapes.add_textbox(Inches(mid - lw / 2), Inches(y - 0.34),
                                       Inches(lw), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.name = BODY
        r.font.color.rgb = color
    return c


def diagram_four_parts(prs):
    s = title_only(prs, "Inside one machine",
                   "Four parts. Every machine you will run research code on has all four.")
    _box(s, 0.40, 1.60, 2.05, 1.30, "CPU", "the processor chip")
    _box(s, 2.70, 1.60, 2.05, 1.30, "Cores", "independent workers inside it")
    _box(s, 5.00, 1.60, 2.05, 1.30, "RAM", "holds what the CPU is using now")
    _box(s, 7.30, 1.60, 2.05, 1.30, "Storage", "where files live at rest")
    box = s.shapes.add_textbox(Inches(0.40), Inches(3.20), Inches(8.95), Inches(1.4))
    write(box.text_frame, [
        ("I/O — moving data from storage into RAM and writing results back.", True),
        "Almost always the slowest thing your script does.",
        "",
        "Your script is just the sequence of steps the CPU follows to turn one into the other.",
    ], size=13)
    return s


def diagram_data_moves(prs):
    s = title_only(prs, "How your data moves",
                   "Storage → RAM → CPU → back to storage.")
    _box(s, 0.45, 1.95, 2.30, 1.15, "Storage", "large, slow")
    _box(s, 3.60, 1.95, 2.30, 1.15, "RAM", "fast, limited")
    _box(s, 6.75, 1.95, 2.30, 1.15, "CPU", "cores do the work")
    _arrow(s, 2.85, 2.52, 3.50, "read — milliseconds")
    _arrow(s, 6.00, 2.52, 6.65, "nanoseconds")
    _arrow(s, 0.45, 3.55, 9.05, "write back — milliseconds again", color=MUTED)
    box = s.shapes.add_textbox(Inches(0.45), Inches(3.85), Inches(8.6), Inches(0.9))
    write(box.text_frame, [
        "The two slow legs are both at the storage end. Once the data is close to the "
        "processor, the work itself is quick.",
    ], size=13)
    return s


def diagram_fan_out(prs):
    s = title_only(prs, "One script in, many tasks out",
                   "You submit once. Slurm runs the same script many times, "
                   "handing each copy a different number.")
    _box(s, 0.45, 2.30, 2.15, 1.00, "one script", "submitted once", fill=TBL_BAND)
    for i, (lbl, y) in enumerate([("task 1", 1.45), ("task 2", 2.35),
                                  ("task 3", 3.25), ("task N", 4.15)]):
        _box(s, 5.20, y, 2.20, 0.68, lbl,
             "SLURM_ARRAY_TASK_ID", label_size=12)
        _arrow(s, 2.75, 2.80, 5.05)
    box = s.shapes.add_textbox(Inches(0.45), Inches(3.70), Inches(4.30), Inches(1.0))
    write(box.text_frame, [
        ("Nothing decides which filing each task takes.", True),
        "That mapping is yours to write.",
    ], size=12.5)
    return s


def diagram_pipeline(prs):
    s = title_only(prs, "Today: one pipeline, five steps",
                   "Part 1 is steps 1 to 3. Part 2 is steps 4 and 5 — and it runs "
                   "on the numbers you write down this morning.")
    steps = ["profile", "submit to Slurm", "read logs", "document", "scale (arrays)"]
    for i, label in enumerate(steps):
        x = 0.40 + i * 1.83
        fill = TBL_BAND if i < 3 else WHITE
        _box(s, x, 2.05, 1.62, 1.05, f"{i + 1}", label, fill=fill, label_size=17)
        if i < len(steps) - 1:
            _arrow(s, x + 1.66, 2.58, x + 1.79)
    box = s.shapes.add_textbox(Inches(0.40), Inches(3.45), Inches(8.95), Inches(0.9))
    write(box.text_frame, [
        ("Everything you type into a job script today comes from a number you measured.",
         True),
        "There is nothing to learn from step four if you skipped step one.",
    ], size=13)
    return s


# ------------------------------------------------------------------ deck ----

WEB = "darc.stanford.edu/class"


def build(template, out):
    prs = Presentation(template)
    strip_slides(prs)

    # ---- front matter
    title_slide(prs, "Day 2 · The Cluster", "September 18, 2026",
                "Alex Storer, Jeff Ott, Natalya Rapstine",
                "Data, Analytics, and Research Computing – GSB Research Hub")

    title_body(prs, "Agenda", [
        ("Part 1 — The cluster, and why it has a queue", True),
        "What is inside any computer, and what running a script actually means",
        "How data moves: disk, RAM, CPU",
        "What profiling is, and why we do it",
        "The queue: interactive vs. batch, partitions, and reading a failed job",
        "",
        ("Part 2 — Doing many things at once", True),
        "What parallelization is, and when it does nothing at all",
        "The three shapes: more cores, more jobs, or both",
        "What a Slurm array is",
        "A look at the GPUs",
        "",
        ("A little lecture, a lot of lab. Two work blocks, self-paced.", True),
    ], size=14.5)

    # ================================================================ PART 1
    section(prs, "Part 1 · The cluster, and why it has a queue",
            "9:00–9:20 · concepts only, laptops closed")

    section(prs, "What is inside any computer?")
    diagram_four_parts(prs)

    table_slide(prs, "The same parts, in different amounts",
                ["What you get", "Your laptop", "One Yen node", "The cloud"],
                [["CPU cores", "a handful", "up to 256 on yen1", "as many as you rent"],
                 ["RAM", "single-digit GB", "250 GB to 3 TB", "as much as you rent"],
                 ["Storage", "your own disk", "shared, ~1 PB, every node sees it",
                  "rented volumes"]],
                sub="What changes between environments is how much of each part you get.",
                col_widths=[1.5, 1.5, 2.0, 1.5])

    table_slide(prs, "And what it is like to use",
                ["", "Your laptop", "One Yen node", "The cloud"],
                [["Who else is on it", "nobody", "everyone — it is contended", "nobody"],
                 ["Close the lid", "the work stops", "it keeps running", "it keeps running"],
                 ["What it costs you", "nothing extra", "nothing extra",
                  "per hour, used or idle"]],
                col_widths=[1.7, 1.4, 1.9, 1.5])

    takeaway(prs, "For most GSB research, the Yens are the right answer.", [
        "No per-hour cost, and the data stays on Stanford-managed infrastructure. "
        "The cloud earns its keep when you need something the Yens do not have — "
        "and you pay for it by the hour whether you are using it or not.",
    ])

    section(prs, "Running a script")

    title_body(prs, "What running a script actually means", [
        ("1 · Load from disk", True),
        "Python reads your script and your data files from storage.",
        ("2 · Into RAM", True),
        "The data lands in memory, where the CPU can reach it quickly.",
        ("3 · The CPU works", True),
        "Cores execute the steps in your script against what is in RAM.",
        ("4 · Save to disk", True),
        "Results are written back to storage, so they survive the run.",
        "",
        ("Watch the first leg. Getting data from disk into RAM is the slow one.", True),
    ], size=14)

    diagram_data_moves(prs)

    takeaway(prs, "Disk is about a million times farther away than RAM.", [
        "Your CPU reaches data in RAM in nanoseconds; a disk read takes milliseconds. "
        "If your dataset does not fit in RAM all at once, your script keeps going back "
        "to disk in the middle of the computation — and that, far more often than a slow "
        "CPU, is what makes a job crawl.",
        ("Which is why knowing how much RAM your script needs matters.", True),
    ])

    section(prs, "Profiling")

    title_body(prs, "What is profiling?", [
        ("Measuring what your code actually consumes as it runs — "
         "instead of guessing at it.", True),
        "",
        ("Time — how long you actually waited for it to finish.", True),
        ("Cores — how many workers it used at once.", True),
        ("RAM — how much memory it held while it ran.", True),
        "",
        "Three numbers. They are the ones the scheduler will ask you for — and the ones "
        "you cannot answer by reading the code.",
        "",
        "This works whether you wrote the script or someone handed it to you. In the "
        "work block you will do it to a script you have never seen.",
    ], size=14)

    title_body(prs, "Why do we profile?", [
        ("Because on a cluster you have to declare what you need "
         "before the job runs.", True),
        "",
        ("Ask too little — the job dies.", True),
        "It hits the ceiling you declared and gets killed. You queued, you waited, "
        "and you have nothing.",
        ("Ask too much — you wait, and you waste.", True),
        "Harder to schedule, so you sit in the queue longer — and you hold capacity "
        "nobody else can use.",
        ("Measure first — right-sized.", True),
        "You ask for roughly what you used. The job starts sooner, finishes, and leaves "
        "the rest of the cluster alone.",
    ], size=14)

    table_slide(prs, "What you measure, and what tells you",
                ["What you want to know", "What tells you", "What you read off it"],
                [["How long it took", "`time`", "real — the clock time you waited"],
                 ["How many cores it used", "`userload`",
                  "Cores — your whole footprint on the node"],
                 ["Which processes, how hard", "`htop`",
                  "CPU% per process; over 100% means several cores"],
                 ["How much RAM it held", "`htop`", "RES — the real memory in use"]],
                note=[("Two terminals, same node.", True),
                      "One runs the script, the other watches it. Both tools only see the "
                      "machine they are on, so the second terminal has to be on the same Yen.",
                      "Small is not zero — a few MB on a 1 TB node rounds to 0.0%. Write "
                      "down a small round number rather than nothing at all."],
                col_widths=[2.0, 1.2, 2.6])

    title_body(prs, "Two numbers tell you the shape of the job", [
        ("Compare the clock time against the CPU time and the script "
         "tells you what kind of thing it is.", True),
        "",
        ("user ≈ real  →  Serial", True),
        "One core at a time. More cores would not help it.",
        ("user > real  →  Parallel", True),
        "Several cores at once — CPU time adds up faster than the clock runs.",
        ("real ≫ user  →  Waiting, or I/O-bound", True),
        "Barely touching the CPU. It is waiting on something: a disk, a network, an API.",
        "",
        "You will meet two of these before 10:30. The mystery script is CPU-bound; the "
        "extraction you wrote yesterday spends almost all its time waiting on the API.",
    ], size=13.5)

    section(prs, "The queue")

    title_body(prs, "What Slurm is", [
        ("The thing that decides who gets which cores, and when.", True),
        "",
        "Seventeen nodes, one shared file system, and everyone at the GSB doing research "
        "wants them at once. Slurm is the queue that settles it.",
        "",
        ("The scarce resource is not a machine. It is core-minutes.", True),
        "You are not asking for a computer — you are asking for a slice of cores for a "
        "span of time, and you have to say how much of each.",
        "",
        ("Nobody is watching you.", True),
        "Once a job is submitted it runs whether you are at your desk or not. That is the "
        "point of handing it over.",
    ], size=14)

    title_body(prs, "Interactive, or handed to Slurm", [
        ("The question is: are you sitting there while it runs?", True),
        "",
        ("Interactive — on a Yen, right now", True),
        "You type, it answers. Good for finding out what your code does, and for anything "
        "you need to watch. You are sharing the node with everyone else logged in.",
        ("Batch — submitted to the queue", True),
        "You declare what you need, hand it over, and walk away. Dedicated cores, and the "
        "job survives you closing your laptop.",
        "",
        ("So the rule is: explore interactively, execute in batch.", True),
    ], size=14)

    title_body(prs, "Partitions and limits", [
        ("A partition is a named pool of machines, each with its own rules.", True),
        "",
        "Different queues have different caps — on how long a job may run, how many cores "
        "it may take, how many tasks an array may hold.",
        "",
        ("The caps are a fairness mechanism, not bureaucracy.", True),
        "They are what stops one person's thousand-task array from holding the cluster "
        "shut for everybody else.",
        "",
        "Which partition you pick therefore decides what you are allowed to ask for. "
        "Check before you assume.",
    ], size=14)

    title_body(prs, "Debugging a job that failed", [
        ("A job that vanished is not a job that worked.", True),
        "",
        "When a batch job fails there is nobody watching and no traceback on your screen. "
        "The cluster writes down what went wrong and moves on — your job is to go and read it.",
        "",
        ("Everything you need is in the .out and .err files.", True),
        "If you do not know where those landed, you cannot debug anything.",
        "",
        ("Why we practise it deliberately:", True),
        "at a hundred tasks, failure stops being exceptional. Reading a failed job's .err "
        "is the first thing you will actually need, and the most often.",
    ], size=13.5)

    diagram_pipeline(prs)

    # ---- Part 1 lab
    section(prs, "Lab Time.", "Part 1 · 9:20–10:30 · work at your own pace")

    title_body(prs, "Your Goals — Part 1", [
        "Profile a script you have never seen — two terminals, same node",
        "Profile the real extraction on 10 filings, and name it serial, parallel or I/O-bound",
        "Write the Resource Profile into your README — the numbers everything else runs on",
        "Peek at the queue, and tell a running job from a pending one",
        "Write and submit a Slurm job, then find its output and error files",
        "Debug a job that failed — the one skill you will need first, and often",
        "",
        (f"Find the lab materials at: {WEB}", True),
        "",
        ("New to the terminal?", True),
        "Getting through the profiling exercises is a good morning. Ask early — that is "
        "what we are here for.",
        ("Done already?", True),
        "Check whether anyone at your table is stuck before you start the bonus material.",
    ], size=13)

    two_col(prs, "Checkpoint — Part 1",
            "Everyone should have reached", [
                "Profiled the mystery script (two terminals)",
                "Profiled the batch script on 10 filings",
                "Resource Profile written into the README ← protect this one",
                "Peeked at the queue: told R from PD",
                "Written and submitted a Slurm job",
                "Debugged a failed job from its .err ← matters most",
                "",
                "Pages: /day2/profiling/ · /day2/slurm-scheduler/ · /day2/slurm-job/",
            ],
            "Extra, if you are ahead", [
                "Watch a job on its node while it runs",
                "An interactive job with srun",
                "fix_me_2 and fix_me_3",
                "Fix the broken one-file script",
                "Chain two jobs with a dependency",
                "The dev partition · vectorization · core counts",
                "",
                "All bonus. None of it is the floor.",
            ],
            note=["If you have not written the Resource Profile into your README, do that "
                  "before anything in the right-hand column — Part 2 runs on those numbers."])

    # ================================================================ PART 2
    section(prs, "Part 2 · Doing many things at once",
            "10:30–10:50 · concepts only, laptops closed again")

    title_body(prs, "The script you just ran", [
        ("You already have a hundred separate jobs.", True),
        "",
        "No filing needs any other filing to be finished first. Running them one after "
        "another was a decision in your code — not something the problem asked for.",
        "",
        "So the rest of this is about undoing that decision.",
    ], size=15)

    section(prs, "Parallelization", "What it is, when it helps, and when it does nothing.")

    title_body(prs, "What it is, and is not", [
        ("Running many copies of the same work at the same time.", True),
        "",
        ("It is not making one run faster.", True),
        "Work that will not come apart still gets plenty from the cluster — a node to "
        "itself, far more memory, nobody competing for your cores, and the freedom to "
        "walk away. It just will not get this.",
        "",
        ("The test is independence.", True),
        "Could you hand each piece to a different person and never let them talk?",
    ], size=14.5)

    title_body(prs, "Why their job comes apart", [
        ("Inside one filing, a chain. Across filings, nothing.", True),
        "",
        "Inside a single filing the steps are strictly ordered: read from disk, send to "
        "the API, validate the reply, write the result. Four cores on one filing leaves "
        "three of them idle.",
        "",
        "Across filings there is no ordering at all. Filing 2 does not care whether "
        "filing 1 is done. Each one is self-contained.",
        "",
        ("So the unit you parallelize is the filing, not the step.", True),
    ], size=14)

    takeaway(prs, "The part that will not split sets the floor.", [
        "However many workers you add, you can only ever shrink the part that does split. "
        "Asking for more cores past that point buys you nothing but a longer wait in the "
        "queue.",
    ], kicker="THE CEILING")

    title_body(prs, "Three from your own research", [
        ("Which of these come apart? Discussion — still no keyboards.", True),
        "",
        ("1 · Sum every value in an array. Then sum an expensive function of every value.",
         True),
        "Both split. Only one is worth it — adding numbers is so cheap that handing "
        "pieces out and collecting them back costs more than the work.",
        ("2 · Check a key is unique before you merge on it.", True),
        "It splits, but not cleanly. Each worker checks its own chunk, but a duplicate can "
        "sit across a boundary — so you still need one step that looks at all of it.",
        ("3 · A hundred pages scraped, every one appending a row to the same CSV.", True),
        "The scraping is independent. The writing is not. A hundred workers appending to "
        "one file interleave, and you get corrupted rows.",
    ], size=12.5)

    takeaway(prs, "Independent compute is easy. Shared state is where it breaks.", [
        ("Write to separate places. Combine afterwards.", True),
        "That single habit is what makes the difference between an array that works and "
        "one that quietly corrupts its own output.",
    ], kicker="THE PATTERN UNDER ALL THREE")

    section(prs, "Arrays, and the alternatives",
            "One way to spread work out. Not the only one.")

    table_slide(prs, "There are three shapes",
                ["Shape", "What it is", "What it costs you"],
                [["One job, many cores",
                  "You split the work in your own code",
                  "Capped at one machine"],
                 ["Many jobs, one core",
                  "An array — Slurm hands out many small jobs",
                  "Scales past one machine; one failure costs one task"],
                 ["Both at once",
                  "Several jobs, each using several cores",
                  "Most throughput, biggest request, longest wait to start"]],
                sub="Where you started this morning was one job, one core: a loop.",
                col_widths=[1.6, 2.2, 2.6])

    table_slide(prs, "What each shape actually costs",
                ["Eight filings, five seconds each", "One core per job",
                 "Many cores per job"],
                [["One job", "40s — this morning's loop, no speedup at all",
                  "20s — capped at one machine"],
                 ["Many jobs", "20s — an array; scales past one machine",
                  "10s — most throughput, biggest ask"]],
                col_widths=[1.6, 2.4, 2.4])

    title_body(prs, "So what is an array", [
        ("One script, submitted once, run many times over.", True),
        "",
        "Every task runs the identical script. Exactly one thing differs between them: a "
        "number saying which task this is.",
        "",
        ("Which means nothing decides which filing each task takes.", True),
        "That mapping is yours to write. Slurm hands you a number; turning that number "
        "into a filing is your code's job.",
    ], size=15)

    diagram_fan_out(prs)

    title_body(prs, "Two things that will bite you", [
        ("You choose where the numbering starts. Then you live with it.", True),
        "Start at zero and the task number is already a list position. Start at one and "
        "every lookup needs a − 1. Forget, and nothing complains — you silently skip the "
        "first one and run off the end.",
        "",
        ("An array is not unlimited.", True),
        "On the normal partition the task numbers run to 511, and the ceiling differs by "
        "partition — so check before you assume. A job with more tasks than that has to "
        "be split into several arrays.",
    ], size=14)

    takeaway(prs, "Estimate, request, run, check.", [
        "Write the estimate down before you submit anything. Turn it into a request. Run "
        "it. Then compare what it really used against what you guessed.",
        ("There is nothing to learn from step four if you skipped step one.", True),
    ], kicker="THE HABIT THAT OUTLASTS TODAY")

    section(prs, "A look at the GPUs", "Preview only. The bonus page has the rest.")

    table_slide(prs, "Fourteen GPUs, in three tiers",
                ["GPU", "Memory on the card"],
                [["A30", "24 GB"], ["A40", "48 GB"], ["H200", "141 GB"]],
                note=["Fourteen in total, against hundreds of CPU cores per node. "
                      "The scarcest thing on the cluster."],
                col_widths=[1.0, 2.0])

    takeaway(prs, "A model that does not fit in GPU memory does not run slowly. "
                  "It does not run.", [
        "Most limits make things slower. This one is a cliff: the weights fit on the card, "
        "or the job will not start at all. VRAM is what decides which models you can load.",
    ], kicker="WHAT DECIDES WHAT YOU CAN RUN")

    title_body(prs, "Would a GPU have helped this morning?", [
        ("You already measured the answer. Remember what your two clocks said.", True),
        "",
        ("No — and not by a little.", True),
        "A GPU makes arithmetic fast. Your job barely does any: it asks for a filing, "
        "waits, and writes the answer down. real ≫ user said so.",
        "",
        ("Where they do earn their keep:", True),
        "Training or fine-tuning a model · large matrix operations · running an LLM's "
        "weights yourself instead of calling someone else's API.",
    ], size=14)

    # ---- Part 2 lab
    section(prs, "Lab Time.", "Part 2 · 10:50–11:50 · hands on keyboards")

    title_body(prs, "Your Goals — Part 2", [
        "Run a hundred filings as an array, and watch it fan out",
        "Make your tasks safe to run again — four lines, and reruns stop costing money",
        "The capstone: estimate first, submit, then check against sacct",
        "",
        (f"Find the lab materials at: {WEB}", True),
        "",
        ("Start from the array script that is already in the repo.", True),
        "slurm/hello_array.slurm exists so that everyone sees fan-out work with their own "
        "eyes. Writing both array files from scratch is bonus, not the floor.",
        "",
        ("By 11:50 the estimate matters more than the run.", True),
        "Committing to a number before you run it is the transferable discipline of the "
        "whole course.",
    ], size=13)

    two_col(prs, "Checkpoint — Part 2",
            "Everyone should have reached", [
                "An array running from slurm/hello_array.slurm",
                "Tasks made rerun-safe (idempotent)",
                "Capstone: estimate written down first",
                "Capstone: submitted, then compared against sacct",
                "",
                "Pages: /day2/job-arrays/ · /day2/capstone/",
            ],
            "Extra, if you are ahead", [
                "Write both array files from scratch",
                "Merge the results into one CSV",
                "Distill a Slurm skill for Claude — /day2/slurm-with-claude/",
                "GPUs — /day2/gpus/",
                "Run all 992 filings",
                "",
                "All bonus. Nobody is behind for skipping these.",
            ],
            note=["Experts: do not just leave. The bonus pages are where the "
                  "Yen-specific tooling lives — quotas, permissions, user limits."])

    prs.save(out)
    return prs


if __name__ == "__main__":
    template = sys.argv[1] if len(sys.argv) > 1 else "day1-template.pptx"
    out = sys.argv[2] if len(sys.argv) > 2 else "day2-combined.pptx"
    deck = build(template, out)
    print(f"wrote {out}: {len(deck.slides)} slides")
