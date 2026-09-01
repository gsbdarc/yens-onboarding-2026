#!/usr/bin/env python3
"""Build the Day 2 lecture deck (v2) -- concept-led, with the website's own graphics.

Reuses the Day 1 theme (masters, layouts, fonts, GSB branding) by building on a
.pptx export of the Day 1 Google Slides deck, and the slide builders already
written in build_day2_deck.py.

Graphics come from render_website_graphics.py, which renders the docs pages'
inline SVGs with headless Chrome -- animated ones as GIFs, which Google Slides
plays in present mode.

    python build_day2_deck_v2.py <day1-template.pptx> <graphics-dir> <out.pptx>
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from PIL import Image  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

from build_day2_deck import (  # noqa: E402
    ACCENT, BODY, HEAD, MUTED,
    L_TITLE_ONLY,
    drop_unused, layout, section, strip_slides, table_slide, takeaway,
    title_body, title_only, title_slide, two_col, write,
)
from pptx import Presentation  # noqa: E402

GFX = None  # set in build()
WEB = "darc.stanford.edu/class"

# content box for pictures
BOX_L, BOX_R = 0.34, 9.66
BOX_TOP_PLAIN, BOX_TOP_SUB = 1.15, 1.44
BOX_BOTTOM = 5.02          # leaves room for a caption strip
CAP_TOP = 5.06


def picture_slide(prs, title, image, sub=None, caption=None):
    """TITLE_ONLY plus a website graphic, scaled to fit and centred."""
    s = title_only(prs, title, sub=sub)
    path = GFX / image
    if not path.exists():
        raise SystemExit(f"missing graphic: {path}")

    with Image.open(path) as im:
        iw, ih = im.size
    top = BOX_TOP_SUB if sub else BOX_TOP_PLAIN
    avail_w = BOX_R - BOX_L
    avail_h = (BOX_BOTTOM if caption else BOX_BOTTOM + 0.35) - top
    scale = min(avail_w / iw, avail_h / ih)
    w, h = iw * scale, ih * scale
    left = BOX_L + (avail_w - w) / 2
    s.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))

    if caption:
        box = s.shapes.add_textbox(Inches(BOX_L), Inches(CAP_TOP), Inches(avail_w), Inches(0.5))
        write(box.text_frame, caption, size=12, color=MUTED, space_after=2)
    return s


def build(template, gfxdir, out):
    global GFX
    GFX = Path(gfxdir)
    prs = Presentation(template)
    strip_slides(prs)

    # ------------------------------------------------------------ front matter
    title_slide(prs, "Day 2 · The Cluster", "September 18, 2026",
                "Alex Storer, Jeff Ott, Natalya Rapstine",
                "Data, Analytics, and Research Computing – GSB Research Hub")

    # Day 1's agenda is short noun phrases, not sentences. Match it.
    title_body(prs, "Agenda", [
        ("Part 1", True),
        "What is inside any computer",
        "Running a script",
        "Research computing resources",
        "Profiling",
        "Why profiling matters",
        "Slurm",
        "",
        ("Part 2", True),
        "Parallelization",
        "Job arrays",
        "GPUs",
        "Where to get help",
    ], size=15)

    # ==================================================================== PART 1
    section(prs, "Part 1 · Where your code runs", "9:00–9:20 · concepts only")

    # ---- what is inside any computer
    section(prs, "What is inside any computer?")

    picture_slide(prs, "Your laptop, a Yen, the cloud", "cluster-shape.png",
                  sub="Same parts everywhere. What changes is how much, and who else wants it.",
                  caption=[("You write code on your laptop. It runs on a machine somewhere else.",
                            True)])

    table_slide(prs, "Same parts, different amounts",
                ["", "Your laptop", "One Yen node", "The cloud"],
                [["CPU cores", "a handful", "up to 256 on yen1", "as many as you rent"],
                 ["RAM", "single-digit GB", "250 GB – 3 TB", "as much as you rent"],
                 ["Storage", "your own disk", "shared, ~1 PB", "rented volumes"],
                 ["Who else is on it", "nobody", "everyone", "nobody"],
                 ["Close the lid", "work stops", "keeps running", "keeps running"]],
                col_widths=[1.6, 1.4, 1.8, 1.6])

    title_body(prs, "CPU and cores", [
        ("The CPU is the processor chip. It does the work.", True),
        "",
        ("A core is one independent worker inside it.", True),
        "Each core runs instructions on its own. That independence is the whole reason "
        "parallel work is possible at all.",
        "",
        "One core can only do one thing at a time. Eight cores can do eight — but only if "
        "your work comes apart into eight pieces.",
    ], size=15)

    title_body(prs, "RAM, storage, and I/O", [
        ("RAM — the desk.", True),
        "Holds the data the CPU is actively using. Fast, and limited. Run out and the job dies.",
        "",
        ("Storage — the filing cabinet.", True),
        "Where your files live when nothing is running. Large, and slow to reach.",
        "",
        ("I/O — walking between them.", True),
        "Moving data from storage into RAM, and writing results back. Almost always the "
        "slowest thing your script does.",
    ], size=14.5)

    # ---- running a script
    section(prs, "Running a script")

    title_body(prs, "What running a script actually means", [
        ("Four things happen, in this order, every time.", True),
        "",
        ("1 · Load from disk", True),
        "Python reads your script and your data files from storage.",
        ("2 · Into RAM", True),
        "The data lands in memory, where the CPU can reach it.",
        ("3 · The CPU works", True),
        "Cores execute your steps against what is in RAM.",
        ("4 · Save to disk", True),
        "Results are written back, so they survive the run.",
    ], size=14)

    picture_slide(prs, "How your data moves", "data-moves.gif",
                  caption=[("Watch the first leg.", True),
                           "Getting data from disk into RAM is the slow one. Once it is close "
                           "to the processor, the work itself is quick."])

    takeaway(prs, "Disk is about a million times farther away than RAM.", [
        "The CPU reaches RAM in nanoseconds; a disk read takes milliseconds. If your data "
        "does not fit in RAM all at once, the script keeps going back to disk mid-computation "
        "— and that, far more often than a slow CPU, is what makes a job crawl.",
    ])

    # ---- resources
    section(prs, "Research computing resources")

    title_body(prs, "You have a running script. What does it need?", [
        ("Three resources account for almost every decision you will make.", True),
        "",
        ("Compute time — how long it runs.", True),
        ("Cores — how many workers it uses at once.", True),
        ("RAM — how much memory it holds while running.", True),
        "",
        "There are others — disk space, network, GPU memory — but these three are the ones "
        "you will declare, argue about, and get wrong.",
    ], size=15)

    title_body(prs, "How much of each does your script use?", [
        ("That is the question. And you cannot answer it by reading the code.", True),
        "",
        "Not from the algorithm, not from the file size, not from how long it took on your "
        "laptop last Tuesday.",
        "",
        ("You have to measure it while it runs.", True),
    ], size=16)

    # ---- profiling
    section(prs, "Profiling")

    title_body(prs, "Measuring those three things is profiling", [
        ("Profiling is measuring what your code actually consumes as it runs — "
         "instead of guessing.", True),
        "",
        ("Memory profiling", True),
        "Measuring how much RAM the job holds.",
        ("CPU profiling", True),
        "Measuring how many cores it uses, and for how long.",
        ("Profiling", True),
        "Usually just means all three at once. That is what we will do.",
        "",
        "It works whether you wrote the script or inherited it. In the work block you will "
        "profile one you have never seen.",
    ], size=13.5)

    picture_slide(prs, "Two terminals, one node", "two-terminals.png",
                  sub="One runs the script. The other watches it.",
                  caption=[("Both tools only see the machine they are on.", True),
                           "So the second terminal has to be on the same Yen as the first."])

    # ---- why profile
    section(prs, "Why profiling matters")

    title_body(prs, "Because resources are finite — and shared", [
        ("Your laptop", True),
        "All yours. Nobody else is competing for it. Also: small, and it sleeps.",
        "",
        ("The cloud", True),
        "All yours too — for rent. You pay by the hour whether you use it or not.",
        "",
        ("The Yens", True),
        "\"Free\" to you, and **shared** with everyone else doing research at the GSB. "
        "Which makes being a good citizen an actual technical skill, not a nicety.",
    ], size=14.5)

    title_body(prs, "What being a good citizen means in practice", [
        ("It means you know what your code needs before you ask for it.", True),
        "",
        "Is this a 2-hour job or a 2-week job?",
        "Does it want 10 cores or 200?",
        "Can the problem be broken into parts and run in parallel to finish sooner?",
        "Are you actually using what you asked for — or holding it idle?",
        "",
        ("Every one of those is a measurement, not an opinion.", True),
    ], size=15)

    title_body(prs, "Getting it wrong costs, in both directions", [
        ("Ask for too much", True),
        "You wait longer in the queue, because a big request is harder to schedule. And once "
        "it starts you hold cores and memory you never touch — which nobody else can use.",
        "",
        ("Ask for too little", True),
        "The job hits the ceiling you declared and gets killed partway through. You queued, "
        "you waited, and you have nothing to show.",
        "",
        ("Measure first, and you ask for roughly what you use.", True),
    ], size=14.5)

    picture_slide(prs, "Three ways to size a request", "size-a-request.png",
                  caption=[("Only the third one is doing your research a favour.", True)])

    # ---- two situations
    section(prs, "Two situations", "The same inherited script, two different mornings.")

    title_body(prs, "Situation 1 — the script you inherited", [
        ("Your PI hands you a script from a PhD student who graduated. "
         "\"Run this, then extend the analysis.\"", True),
        "",
        "1 · Copy it to the Yens",
        "2 · Profile it on an interactive Yen",
        "3 · It measures 10 GB of RAM, 25 cores, 5 hours",
        "4 · You run it right there on the interactive Yen — and sit watching the output for "
        "five hours, hoping your connection does not drop",
        "",
        ("Step 4 is the problem. Nothing about it is wrong except everything.", True),
    ], size=14)

    title_body(prs, "Situation 2 — same script, worse morning", [
        ("Same script, same instruction.", True),
        "",
        "1 · Copy it to the Yens",
        "2 · Profile it on an interactive Yen — and partway through it gets **killed** for "
        "exceeding the per-user memory limit",
        "3 · You scale the input down to 10% of the data and profile again. This time it "
        "runs: 10 GB of RAM, 1 core, 1 hour",
        "",
        ("But that is a tenth of the data.", True),
        "So what do you actually ask for?",
    ], size=14)

    title_body(prs, "The correct way, both times", [
        ("Stop running it on the interactive Yen. Submit it to yen-slurm.", True),
        "",
        ("Situation 1 — you measured the real thing", True),
        "Ask for 10 GB, 25 cores, 5 hours.",
        "",
        ("Situation 2 — you measured a tenth, so scale the estimate", True),
        "Ask for 100 GB, 1 core, 10 hours.",
        "",
        ("Then close your laptop and walk away.", True),
        "Work on another part of the project. Have lunch. Take a nap. The job does not need "
        "you in the room.",
    ], size=14)

    title_body(prs, "It is subtler than multiplying by ten", [
        ("Not all code scales linearly.", True),
        "Ten times the data is not always ten times the memory, and almost never ten times "
        "the time. Sometimes it is worse than linear.",
        "",
        ("And the estimate is not the only lever.", True),
        "Maybe the right move is to rewrite the script to use more cores and less RAM. Maybe "
        "it is to break the work into independent pieces and run them at the same time.",
        "",
        ("That second idea is Part 2.", True),
    ], size=14.5)

    # ---- slurm
    section(prs, "Slurm")

    title_body(prs, "What Slurm gives you", [
        ("A queue, and then a machine that is yours.", True),
        "",
        ("More resources than interactive.", True),
        "The limits that killed the job in Situation 2 are interactive limits. Batch nodes "
        "are bigger.",
        "",
        ("And they are yours while the job runs.", True),
        "Dedicated cores and memory — nobody else competing, nobody else's job slowing yours.",
        "",
        ("The price: you have to say what you need up front.", True),
        "Which is why the whole morning has been about measuring it.",
    ], size=14)

    picture_slide(prs, "Where a job goes", "slurm-submit.png",
                  sub="You submit from a shared interactive Yen. Slurm finds you a compute node.",
                  caption=[("You are not asking for a computer.", True),
                           "You are asking for a slice of cores, for a span of time."])

    picture_slide(prs, "The lifecycle of a job", "job-lifecycle.gif",
                  caption=[("PD means queued and waiting for a node. R means running.", True),
                           "When it finishes, what it printed is waiting for you in logs/."])

    picture_slide(prs, "Where this morning fits", "day2-map.png",
                  caption=[("Part 1 is steps 1 to 3. Part 2 is 4 and 5 — and it runs on the "
                            "numbers you write down this morning.", True)])

    # ---- part 1 lab
    section(prs, "Lab Time.", "Part 1 · 9:20–10:30 · work at your own pace")

    title_body(prs, "Your Goals — Part 1", [
        "Profile a script you have never seen",
        "Write the resources down — time, cores, RAM",
        "Submit a Slurm job",
        "Read its log",
        "",
        (f"Find the lab materials at: {WEB}", True),
        "",
        ("New to the terminal?", True),
        "Getting through the profiling exercises is a good morning. Ask early — that is what "
        "we are here for.",
        ("Done already?", True),
        "Check whether anyone at your table is stuck before you start the bonus material.",
    ], size=14)

    two_col(prs, "Checkpoint — Part 1",
            "Everyone should have reached", [
                "Profiled the mystery script (two terminals)",
                "Profiled the batch script on 10 filings",
                "Resource Profile written into the README ← protect this one",
                "Peeked at the queue: told R from PD",
                "Written and submitted a Slurm job",
                "Debugged a failed job from its .err ← matters most",
                "",
                "/day2/profiling/ · /day2/slurm-scheduler/ · /day2/slurm-job/",
            ],
            "Extra, if you are ahead", [
                "Watch a job on its node while it runs",
                "An interactive job with srun",
                "fix_me_2 and fix_me_3",
                "Fix the broken one-file script",
                "Chain two jobs with a dependency",
                "The dev partition · vectorization · core counts",
                "",
                "All bonus. None of it is the floor.",
            ],
            note=["Not done the Resource Profile? Do that before anything on the right — "
                  "Part 2 runs on those numbers."])

    # ==================================================================== PART 2
    section(prs, "Part 2 · Doing many things at once", "10:30–10:50 · concepts only")

    title_body(prs, "The script you just ran", [
        ("You already have a hundred separate jobs.", True),
        "",
        "No filing needs any other filing to be finished first. Running them one after "
        "another was a decision in your code — not something the problem asked for.",
        "",
        ("The rest of this is about undoing that decision.", True),
    ], size=15.5)

    # ---- parallelization
    section(prs, "Parallelization", "What it is, when it helps, and when it does nothing.")

    title_body(prs, "What it is, and is not", [
        ("Running many copies of the same work at the same time.", True),
        "",
        ("It is not making one run faster.", True),
        "Work that will not come apart still gets plenty from the cluster — a node to itself, "
        "far more memory, nobody competing for your cores, and the freedom to walk away. It "
        "just will not get this.",
        "",
        ("The test is independence.", True),
        "Could you hand each piece to a different person, and never let them talk?",
    ], size=14.5)

    picture_slide(prs, "One burner, four sandwiches", "kitchen-one.gif",
                  caption=[("The steps of one sandwich are a chain.", True),
                           "Grilling cannot start before assembly. More cooks change nothing."])

    picture_slide(prs, "Four burners, four sandwiches", "kitchen-four.gif",
                  caption=[("The sandwiches are independent.", True),
                           "Four burners finish at t = 4 what one burner finishes at t = 16."])

    title_body(prs, "Your extraction job is four burners", [
        ("Inside one filing, a chain. Across filings, nothing.", True),
        "",
        "Inside a single filing the steps are strictly ordered: read it, send it to the API, "
        "validate the reply, write the result. Four cores on one filing leaves three idle.",
        "",
        "Across filings there is no ordering at all. Filing 2 does not care whether filing 1 "
        "is done.",
        "",
        ("So the unit you parallelize is the filing, not the step.", True),
    ], size=14.5)

    takeaway(prs, "The part that will not split sets the floor.", [
        "However many workers you add, you can only ever shrink the part that does split. "
        "Asking for more cores past that point buys you nothing but a longer wait in the queue.",
    ], kicker="THE CEILING")

    title_body(prs, "Three from your own research", [
        ("Which of these come apart? Discussion — no keyboards.", True),
        "",
        ("1 · Sum every value in an array. Now sum an expensive function of every value.", True),
        "Both split. Only one is worth it — adding numbers is so cheap that handing pieces "
        "out and collecting them back costs more than the work.",
        "",
        ("2 · Check a key is unique before you merge on it.", True),
        "Splits, but not cleanly. Each worker checks its chunk — but a duplicate can sit "
        "across a boundary, so one step still has to see all of it.",
        "",
        ("3 · A hundred pages scraped, each appending a row to the same CSV.", True),
        "The scraping is independent. The writing is not. A hundred writers interleave, and "
        "you get corrupted rows.",
    ], size=12.5)

    takeaway(prs, "Independent compute is easy. Shared state is where it breaks.", [
        ("Write to separate places. Combine afterwards.", True),
        "That one habit is the difference between an array that works and one that quietly "
        "corrupts its own output.",
    ], kicker="THE PATTERN UNDER ALL THREE")

    # ---- shapes
    section(prs, "Shapes of parallelism", "More cores, more jobs, or both.")

    picture_slide(prs, "Where you started: one job, one core", "shape-1job1core.gif",
                  caption=[("This morning's loop. Eight filings, one after another.", True),
                           "The baseline everything else is measured against."])

    picture_slide(prs, "One job, many cores", "shape-1jobNcore.gif",
                  caption=[("The job splits the work across cores itself.", True),
                           "Faster — and capped at one machine, since a job cannot use cores "
                           "on a node it was not given."])

    picture_slide(prs, "Many jobs, one core each", "shape-Njob1core.gif",
                  caption=[("An array. Slurm hands out many small jobs instead of one big one.",
                            True),
                           "Scales past a single machine, and one failure costs one task "
                           "instead of the run."])

    picture_slide(prs, "Many jobs, many cores", "shape-NjobNcore.gif",
                  caption=[("Both at once. Most throughput, biggest request.", True),
                           "And the longest wait before anything starts."])

    table_slide(prs, "What each shape actually costs",
                ["Eight filings, five seconds each", "One core per job", "Many cores per job"],
                [["One job", "40s — a loop, no speedup", "20s — capped at one machine"],
                 ["Many jobs", "20s — an array, scales out", "10s — most throughput"]],
                col_widths=[1.7, 2.3, 2.3])

    # ---- arrays
    section(prs, "Job arrays")

    title_body(prs, "What an array is", [
        ("One script, submitted once, run many times over.", True),
        "",
        "Every task runs the identical script. Exactly one thing differs between them: a "
        "number saying which task this is.",
        "",
        ("Which means nothing decides which filing each task takes.", True),
        "Slurm hands you a number. Turning that number into a filing is your code's job.",
    ], size=15)

    picture_slide(prs, "One script in, many tasks out", "array-fanout.png",
                  caption=[("You submit once.", True),
                           "Slurm runs the same script many times, handing each copy a "
                           "different task ID."])

    title_body(prs, "Two things that will bite you", [
        ("You choose where the numbering starts. Then you live with it.", True),
        "Start at zero and the task number is already a list position. Start at one and every "
        "lookup needs a − 1. Forget, and nothing complains — you silently skip the first one "
        "and run off the end.",
        "",
        ("An array is not unlimited.", True),
        "On the normal partition task numbers run to 511, and the ceiling differs by "
        "partition. More tasks than that has to be split across several arrays.",
    ], size=14)

    title_body(prs, "Make your tasks safe to run again", [
        ("If the output for this task already exists, skip it.", True),
        "",
        "Four lines of code. It turns a rerun from \"pay for all hundred again\" into "
        "\"finish the six that failed\".",
        "",
        "At a hundred tasks, some will fail. That stops being an exception and becomes the "
        "normal case — so the ability to just resubmit the whole array is what saves you.",
    ], size=15)

    takeaway(prs, "Estimate, request, run, check.", [
        "Write the estimate down before you submit anything. Turn it into a request. Run it. "
        "Then compare what it really used against what you guessed.",
        ("There is nothing to learn from step four if you skipped step one.", True),
    ], kicker="THE HABIT THAT OUTLASTS TODAY")

    # ---- gpus
    section(prs, "GPUs", "Preview only. The bonus page has the rest.")

    table_slide(prs, "Fourteen GPUs, in three tiers",
                ["GPU", "Memory on the card"],
                [["A30", "24 GB"], ["A40", "48 GB"], ["H200", "141 GB"]],
                note=["Fourteen in total, against hundreds of CPU cores per node. "
                      "The scarcest thing on the cluster."],
                col_widths=[1.0, 2.0])

    takeaway(prs, "A model that does not fit in GPU memory does not run slowly. "
                  "It does not run.", [
        "Most limits make things slower. This one is a cliff: the weights fit on the card, or "
        "the job will not start. VRAM decides which models you can load at all.",
    ], kicker="WHAT DECIDES WHAT YOU CAN RUN")

    title_body(prs, "Would a GPU have helped this morning?", [
        ("You already measured the answer.", True),
        "",
        ("No — and not by a little.", True),
        "A GPU makes arithmetic fast. Your job barely does any: it asks for a filing, waits, "
        "and writes the answer down. Your two clocks said so.",
        "",
        ("Where they do earn their keep:", True),
        "Training or fine-tuning a model · large matrix operations · running an LLM's weights "
        "yourself instead of calling someone else's API.",
    ], size=14.5)

    # ---- help
    section(prs, "Where to get help")

    title_body(prs, "After today", [
        ("RCpedia — rcpedia.stanford.edu", True),
        "The written documentation for the Yens. Storage, Slurm, software, quotas.",
        "",
        ("Slack — #gsb-yen-users", True),
        "Where Yen users and the DARC team answer questions about the cluster, Slurm, storage "
        "and software. Also where maintenance windows and new hardware get announced.",
        "",
        ("Email — gsb_darcresearch@stanford.edu", True),
        "For anything you would rather not post in a channel, or that needs a direct answer. "
        "Typically one business day.",
        "",
        ("The course site stays up.", True),
    ], size=13.5)

    title_body(prs, "When you are stuck", [
        ("\"My Slurm job keeps failing\"", True),
        "#gsb-yen-users — someone has seen it. Include your error output.",
        ("\"My code works on my laptop but not the Yens\"", True),
        "#gsb-yen-users — again, paste the error.",
        ("\"Is this dataset ok to send to an LLM?\"", True),
        "Email DARC, or ask your IRB coordinator.",
        ("\"I want to run something much bigger\"", True),
        "Email DARC — we can advise on allocations.",
    ], size=14)

    # ---- part 2 lab
    section(prs, "Lab Time.", "Part 2 · 10:50–11:50 · hands on keyboards")

    title_body(prs, "Your Goals — Part 2", [
        "Submit a job array — and watch it fan out",
        "Make your tasks safe to run again",
        "The capstone: estimate first, then check against what it used",
        "",
        (f"Find the lab materials at: {WEB}", True),
        "",
        ("Start from the array script already in the repo.", True),
        "slurm/hello_array.slurm exists so everyone sees fan-out work with their own eyes. "
        "Writing both array files from scratch is bonus, not the floor.",
        "",
        ("By 11:50 the estimate matters more than the run.", True),
    ], size=14)

    two_col(prs, "Checkpoint — Part 2",
            "Everyone should have reached", [
                "An array running from slurm/hello_array.slurm",
                "Tasks made rerun-safe (idempotent)",
                "Capstone: estimate written down first",
                "Capstone: submitted, then compared against sacct",
                "",
                "/day2/job-arrays/ · /day2/capstone/",
            ],
            "Extra, if you are ahead", [
                "Write both array files from scratch",
                "Merge the results into one CSV",
                "Distill a Slurm skill for Claude — /day2/slurm-with-claude/",
                "GPUs — /day2/gpus/",
                "Run all 992 filings",
                "",
                "All bonus. Nobody is behind for skipping these.",
            ],
            note=["Experts: don't just leave. The bonus pages are where the Yen-specific "
                  "tooling lives — quotas, permissions, user limits."])

    prs.save(out)
    return prs


if __name__ == "__main__":
    tpl = sys.argv[1]
    gfx = sys.argv[2]
    out = sys.argv[3] if len(sys.argv) > 3 else "day2-deck-v2.pptx"
    deck = build(tpl, gfx, out)
    print(f"wrote {out}: {len(deck.slides)} slides")
